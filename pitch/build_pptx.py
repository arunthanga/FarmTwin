"""Generate FarmTwin-KSUM.pptx from the pitch content (python-pptx).

Usage:
    pip install python-pptx
    python pitch/build_pptx.py
Produces: pitch/FarmTwin-KSUM.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Brand palette
INK     = RGBColor(0x0D, 0x1B, 0x2A)
WATER   = RGBColor(0x2A, 0xA7, 0xD8)
WATER_D = RGBColor(0x15, 0x77, 0xA6)
LEAF    = RGBColor(0x4C, 0xAF, 0x50)
EARTH   = RGBColor(0xC8, 0xA0, 0x64)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xEA, 0xF6, 0xFB)
GREY    = RGBColor(0x5B, 0x72, 0x86)
DARKTX  = RGBColor(0x1A, 0x2A, 0x38)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(1, x, y, w, h)  # 1 = rectangle
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def textbox(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            space_after=6):
    """runs: list of paragraphs; each paragraph = list of (text, size, color, bold)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        for (text, size, color, bold) in para:
            r = p.add_run(); r.text = text
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
            r.font.name = "Calibri"
    return tb


def line(slide, x1, y1, x2, y2, color, width_pt):
    conn = slide.shapes.add_connector(2, x1, y1, x2, y2)  # 2 = straight
    conn.line.color.rgb = color
    conn.line.width = Pt(width_pt)
    conn.shadow.inherit = False
    return conn


def oval(slide, x, y, w, h, color, label="", lblcolor=WHITE, size=12):
    shp = slide.shapes.add_shape(9, x, y, w, h)  # 9 = oval
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.color.rgb = WHITE; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    if label:
        tf = shp.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = label
        r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = lblcolor
    return shp


def title_bar(slide, title):
    rect(slide, 0, 0, SW, Inches(1.05), INK)
    rect(slide, 0, Inches(1.05), SW, Inches(0.06), WATER)
    textbox(slide, Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.8),
            [[(title, 30, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)


def bullets(slide, items, x=Inches(0.6), y=Inches(1.4), w=Inches(12.1),
            h=Inches(5.6), size=20):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        # it = (text, level, color, bold) or string
        if isinstance(it, str):
            text, level, color, bold = it, 0, DARKTX, False
        else:
            text, level, color, bold = it
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.space_after = Pt(10)
        r = p.add_run(); r.text = ("•  " if level == 0 else "–  ") + text
        r.font.size = Pt(size - level * 2); r.font.color.rgb = color
        r.font.bold = bold; r.font.name = "Calibri"
    return tb


def table(slide, rows, x=Inches(0.6), y=Inches(1.5), w=Inches(12.1), h=Inches(3.2),
          header=True):
    nr, nc = len(rows), len(rows[0])
    gtbl = slide.shapes.add_table(nr, nc, x, y, w, h).table
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = gtbl.cell(ri, ci)
            cell.text = str(val)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(15)
            para.font.name = "Calibri"
            if header and ri == 0:
                para.font.bold = True; para.font.color.rgb = WHITE
                cell.fill.solid(); cell.fill.fore_color.rgb = WATER_D
            else:
                para.font.color.rgb = DARKTX
                cell.fill.solid(); cell.fill.fore_color.rgb = (
                    RGBColor(0xF2, 0xF7, 0xFB) if ri % 2 else WHITE)
    return gtbl


# ---------------------------------------------------------------- Slide 1: Title
s = add_slide(); bg(s, INK)
rect(s, 0, Inches(6.7), SW, Inches(0.8), WATER_D)
textbox(s, Inches(0.7), Inches(2.1), Inches(12), Inches(1.4),
        [[("Farm", 66, WHITE, True), ("Twin", 66, WATER, True)]])
textbox(s, Inches(0.72), Inches(3.5), Inches(11.8), Inches(1.0),
        [[("A digital twin that makes sure every drop of water reaches the "
           "last farmer's roots.", 26, LEAF, True)]])
textbox(s, Inches(0.72), Inches(4.6), Inches(11.8), Inches(1.6),
        [[("Built in Eruthempathy, Chittur (Palakkad) — Kerala's rain-shadow belt.",
           18, LIGHT, False)],
         [("From 20 years of CAE simulation for the world's factories "
           "→ to water for our own village.", 18, LIGHT, False)],
         [("Kerala Startup Mission   ·   github.com/arunthanga/FarmTwin",
           14, RGBColor(0x9F,0xB3,0xC8), False)]])

# ---------------------------------------------------------------- Slide 2: Story
s = add_slide(); bg(s, WHITE); title_bar(s, "This is personal")
bullets(s, [
    ("In my village, Eruthempathy, farmers run out of water.", 0, DARKTX, True),
    ("We get < 1,000 mm rain/yr — one-third of Kerala's average.", 0, DARKTX, False),
    ("Neighbours queue for hours at the regulator for a tanker's worth of drinking water.", 0, DARKTX, False),
    ("The young leave, because the soil can't pay them back.", 0, DARKTX, False),
], y=Inches(1.45), size=22)
rect(s, Inches(0.6), Inches(4.7), Inches(0.08), Inches(1.6), LEAF)
textbox(s, Inches(0.85), Inches(4.7), Inches(11.6), Inches(1.7),
        [[("\u201cI spent 20 years simulating the world's factories — but my own "
           "village couldn't predict whether water would reach a field next week.\u201d",
           22, WATER_D, True)]])

# ------------------------------------------------------------- Slide 3: Problem
s = add_slide(); bg(s, WHITE); title_bar(s, "The problem (the head, after the heart)")
bullets(s, [
    ("Palakkad's east is Kerala's driest belt; water is contested (shared with Tamil Nadu) and badly distributed.", 0, DARKTX, False),
    ("26 blocks in Palakkad rated unsafe on groundwater.", 0, DARKTX, False),
    ("Old Moolathara lift canal: no summer water for 7 years → ~1,000 acres of coconut hit.", 0, DARKTX, False),
    ("A 2-crore clean-up was announced — and nothing happened.", 0, DARKTX, False),
], y=Inches(1.45), size=21)
textbox(s, Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.4),
        [[("Building canals \u2260 delivering water to roots. The last mile is an "
           "information & control problem, not a concrete one.", 22, WATER_D, True)]])

# ------------------------------------------------------------- Slide 4: Why now
s = add_slide(); bg(s, WHITE); title_bar(s, "Why now — money is flowing into our backyard")
bullets(s, [
    ("Moolathara Right Bank Canal Extension — KIIFB, ~262 cr: Kerala's largest community micro-irrigation project, commissioning early 2026.", 0, DARKTX, False),
    ("Serves 3,575 ha (Phase I) & 10,000+ ha (Phase II) — naming Eruthempathy, Vadakarapathy, Kozhinjampara.", 0, DARKTX, False),
    ("CMI is an active, repeating pipeline: ~40.5 cr already tendered across 3 zones.", 0, DARKTX, False),
    ("Farmers get up to 85% subsidy (PDMC + Kerala Samrudhi).", 0, DARKTX, False),
], y=Inches(1.45), size=20)
textbox(s, Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.0),
        [[("Capital is subsidised. The missing piece is the intelligence layer — that's us.",
           22, LEAF, True)]])

# ------------------------------------------------------------ Slide 5: Solution
s = add_slide(); bg(s, WHITE); title_bar(s, "The solution: FarmTwin")
textbox(s, Inches(0.6), Inches(1.35), Inches(12.1), Inches(0.6),
        [[("A self-calibrating digital twin for irrigation — two products, one learning engine.",
           20, DARKTX, True)]])
rect(s, Inches(0.6), Inches(2.1), Inches(5.9), Inches(3.4), RGBColor(0xEF,0xF7,0xF0))
rect(s, Inches(6.9), Inches(2.1), Inches(5.8), Inches(3.4), RGBColor(0xE7,0xF3,0xF9))
textbox(s, Inches(0.85), Inches(2.3), Inches(5.4), Inches(3.1),
        [[("FarmTwin Studio", 22, LEAF, True)],
         [("BEFORE install", 12, GREY, True)],
         [("Survey land → solver optimises the network (pumps, pipes, valves, "
           "drips, sensors, fertigation) → best 2-3 layouts by cost, energy, "
           "uniformity, yield.", 17, DARKTX, False)]])
textbox(s, Inches(7.15), Inches(2.3), Inches(5.3), Inches(3.1),
        [[("FarmTwin Runtime", 22, WATER_D, True)],
         [("AFTER install", 12, GREY, True)],
         [("Solar, wireless edge controllers read sensors + weather + cloud → "
           "decide when/how long to irrigate & fertigate; act on pumps & valves; "
           "reject bad data.", 17, DARKTX, False)]])
textbox(s, Inches(0.6), Inches(5.7), Inches(12.1), Inches(0.8),
        [[("Underneath: FarmTwin Engine — a hydraulic + FAO-56 agronomy solver whose "
           "parameters keep learning from the field.", 16, GREY, True)]], align=PP_ALIGN.CENTER)

# ----------------------------------------------------------- Slide 6: How it works
s = add_slide(); bg(s, WHITE); title_bar(s, "How it works (the engineering moat)")
flow = ("Survey  →  FarmTwin Studio  →  optimal design + BoM + sensor/valve plan\n"
        "                                   ↓\n"
        "                   install (rides on govt / PDMC capex)\n"
        "                                   ↓\n"
        "sensors + weather  →  FarmTwin Runtime (edge)  →  pumps / valves / fertigation\n"
        "                                   ↓\n"
        "        FarmTwin Engine digital twin assimilates reality,\n"
        "        recalibrates parameters  →  better design next time  (learning loop)")
tb = textbox(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(3.4),
             [[(flow, 15, DARKTX, False)]])
textbox(s, Inches(0.6), Inches(5.2), Inches(12.1), Inches(1.4),
        [[("Same discipline as CAE: discretise → solve the physics → compare to reality "
           "→ correct. Most agritech stops at dashboards. We predict, optimise & control "
           "— and get smarter every farm.", 19, DARKTX, False)]])

# ------------------------------------------------------------- Slide 7: Why us
s = add_slide(); bg(s, WHITE); title_bar(s, "Why us (unfair advantages)")
bullets(s, [
    ("20 years of CAE / simulation — meshing, solvers, virtual lines: the rare skill this needs.", 0, DARKTX, False),
    ("A live pilot in the exact project panchayat — 15 acres in Eruthempathy, on the Moolathara RBC ayacut.", 0, DARKTX, False),
    ("Working code already — open hydraulic + FAO-56 engine (GGA solver, components, emitters, agronomy) + a runnable MVP.", 0, DARKTX, False),
    ("Research on our doorstep — IIT Palakkad (water resources, HPC) + Kerala Agricultural University.", 0, DARKTX, False),
], y=Inches(1.6), size=21)

# ----------------------------------------------------- Slide 8: Immediate win
s = add_slide(); bg(s, WHITE); title_bar(s, "The immediate win: Moolathara RBC")
textbox(s, Inches(0.6), Inches(1.3), Inches(12.1), Inches(0.5),
        [[("A government-funded, de-risked beachhead — 2 km from my home.", 19, DARKTX, True)]])
table(s, [
    ["Zone", "Work", "Value", "Tendered"],
    ["Zone-I", "General Civil Work", "₹13.23 cr", "Feb 2026"],
    ["Zone-II", "CMI General Civil Work", "₹13.15 cr", "Feb 2026"],
    ["Zone-III", "CMI DPR Preparation", "₹14.19 cr", "Feb 2026"],
    ["Total (one canal)", "Already tendered", "≈ ₹40.5 cr", ""],
], y=Inches(1.85), h=Inches(2.7))
textbox(s, Inches(0.6), Inches(4.8), Inches(12.1), Inches(1.4),
        [[("Our role: design + digital-twin + performance-monitoring (with a civil "
           "partner). Target the next packages — Phase-II zones & Phase-I O&M — plus a "
           "KSUM-funded pilot on our 15 acres.", 17, DARKTX, False)],
         [("Source: KIIDC e-tenders. Zone I-III closed Feb 2026 — proof of a live, "
           "repeating pipeline.", 12, GREY, False)]])

# ----------------------------------------- Slide 8b: Command-area map (schematic)
s = add_slide(); bg(s, WHITE); title_bar(s, "Where: the Moolathara command area")
textbox(s, Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.4),
        [[("Chittur taluk, Palakkad — schematic, west \u2190\u2192 east, not to scale.",
           14, GREY, False)]])
# Kerala / TN border
line(s, Inches(10.7), Inches(1.7), Inches(10.7), Inches(6.9), EARTH, 2)
textbox(s, Inches(9.4), Inches(1.6), Inches(1.2), Inches(0.3), [[("KERALA", 11, EARTH, True)]], align=PP_ALIGN.RIGHT)
textbox(s, Inches(10.8), Inches(1.6), Inches(1.8), Inches(0.3), [[("TAMIL NADU", 11, EARTH, True)]])
# TN reservoir source
oval(s, Inches(11.0), Inches(2.0), Inches(1.9), Inches(1.0), WATER_D, "Aliyar /\nParambikulam (TN)", LIGHT, 11)
# Chitturpuzha + regulator
line(s, Inches(10.9), Inches(2.6), Inches(8.9), Inches(3.0), WATER, 6)
textbox(s, Inches(9.2), Inches(2.4), Inches(1.6), Inches(0.3), [[("Chitturpuzha", 11, WATER_D, False)]])
reg = rect(s, Inches(8.55), Inches(2.85), Inches(0.45), Inches(0.7), WATER_D)
textbox(s, Inches(7.9), Inches(2.45), Inches(1.8), Inches(0.3), [[("Moolathara Regulator", 12, DARKTX, True)]])
# Left Bank Canal (existing)
line(s, Inches(8.55), Inches(2.95), Inches(6.2), Inches(2.0), RGBColor(0x5B,0x8F,0xB0), 4)
textbox(s, Inches(4.6), Inches(1.75), Inches(2.0), Inches(0.3), [[("Left Bank Canal (existing)", 11, GREY, False)]])
# Right Bank Canal extension (the project)
line(s, Inches(8.4), Inches(3.3), Inches(5.0), Inches(4.4), WATER, 7)
line(s, Inches(5.0), Inches(4.4), Inches(1.9), Inches(5.6), WATER, 7)
textbox(s, Inches(5.6), Inches(3.2), Inches(4.5), Inches(0.3), [[("Phase I: Korayar\u2192Varattayar · 3,575 ha", 12, LEAF, True)]])
textbox(s, Inches(2.4), Inches(5.0), Inches(4.5), Inches(0.3), [[("Phase II: \u2192Velanthavalam · 10,000+ ha", 12, LEAF, True)]])
# Zone stations
oval(s, Inches(6.7), Inches(3.55), Inches(0.4), Inches(0.4), WATER, "I", WHITE, 12)
oval(s, Inches(4.6), Inches(4.2), Inches(0.4), Inches(0.4), WATER, "II", WHITE, 12)
oval(s, Inches(3.0), Inches(4.85), Inches(0.4), Inches(0.4), WATER, "III", WHITE, 12)
# Panchayat command-area blocks
b1 = rect(s, Inches(7.0), Inches(4.3), Inches(2.4), Inches(1.1), RGBColor(0x1B,0x3A,0x24))
textbox(s, Inches(7.15), Inches(4.4), Inches(2.2), Inches(0.8),
        [[("Eruthempathy", 15, RGBColor(0xBF,0xE6,0xC5), True)], [("(our pilot panchayat)", 11, RGBColor(0x8F,0xB8,0x9A), False)]])
b2 = rect(s, Inches(4.5), Inches(5.1), Inches(2.2), Inches(1.0), RGBColor(0x16,0x33,0x1F))
textbox(s, Inches(4.65), Inches(5.3), Inches(2.0), Inches(0.4), [[("Vadakarapathy", 15, RGBColor(0xBF,0xE6,0xC5), True)]])
b3 = rect(s, Inches(1.9), Inches(5.85), Inches(2.3), Inches(0.95), RGBColor(0x14,0x2D,0x1B))
textbox(s, Inches(2.05), Inches(6.05), Inches(2.0), Inches(0.4), [[("Kozhinjampara", 15, RGBColor(0xBF,0xE6,0xC5), True)]])
# Pilot marker
oval(s, Inches(9.05), Inches(4.75), Inches(0.45), Inches(0.45), RGBColor(0xFF,0xD2,0x4A), "", DARKTX)
textbox(s, Inches(9.55), Inches(4.6), Inches(3.0), Inches(0.6),
        [[("\u2605 15-acre pilot", 13, EARTH, True)], [("PWD road · ~2 km from TN border", 11, GREY, False)]])
textbox(s, Inches(0.6), Inches(6.95), Inches(12.1), Inches(0.4),
        [[("Schematic only — illustrative geometry, not a survey map. "
           "Source: KIIDC / KIIFB Moolathara RBC extension notices.", 11, GREY, False)]])

# ------------------------------------------------- Slide 9: Opportunity size
s = add_slide(); bg(s, WHITE); title_bar(s, "Opportunity size: beachhead → world")
table(s, [
    ["Tier", "Scope", "Addressable (our layer)"],
    ["SOM (Yr 1-2)", "Moolathara Phase I design + KSUM pilot + 1st monitoring deal", "₹1-2 cr"],
    ["Beachhead (2-3 yr)", "Moolathara I+II (~13,500 ha) + Chitturpuzha (20,440 ha)", "₹10-25 cr"],
    ["SAM (3-5 yr)", "Kerala KIIFB-CMI + PDMC micro-irrigation", "₹100-300 cr"],
    ["TAM (long)", "India PMKSY (1000s cr/yr); ~70 M ha; global precision irrigation", "very large"],
], y=Inches(1.7), h=Inches(3.0))
textbox(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.0),
        [[("Start with one canal we can touch — end with farmers everywhere.", 22, WATER_D, True)],
         [("Estimates except cited facts; assumptions in the beachhead doc.", 12, GREY, False)]])

# --------------------------------------------------- Slide 10: Business model
s = add_slide(); bg(s, WHITE); title_bar(s, "Business model (subsidy-aligned)")
bullets(s, [
    ("Design & optimisation (Studio): project fee ~3-6k/ha.", 0, DARKTX, False),
    ("Digital-twin + IoT (Runtime): hardware + software — fundable inside the ~85% micro-irrigation subsidy.", 0, DARKTX, False),
    ("Recurring SaaS + O&M / performance monitoring: ~1,000-1,500/ha/yr — sticky, scalable, fixes the 'nobody maintained it' failure.", 0, DARKTX, False),
    ("Phase 2: data-driven agri-fintech (yield-backed credit, insurance).", 0, DARKTX, False),
], y=Inches(1.6), size=21)

# --------------------------------------------------------- Slide 11: GTM
s = add_slide(); bg(s, WHITE); title_bar(s, "Go-to-market")
bullets(s, [
    ("Prove — 15 acres in Eruthempathy + one Moolathara CMI zone (KSUM pilot).", 0, DARKTX, False),
    ("Win — Moolathara design / monitoring with a civil partner.", 0, DARKTX, False),
    ("Expand — Chitturpuzha command area + Palakkad FPOs.", 0, DARKTX, False),
    ("Replicate — KIIFB-CMI statewide (Kerala) & PDMC nationally.", 0, DARKTX, False),
    ("Globalise — software / twin (near-zero marginal cost).", 0, DARKTX, False),
], y=Inches(1.5), size=21)

# ----------------------------------------------------- Slide 12: Traction
s = add_slide(); bg(s, WHITE); title_bar(s, "Traction & assets")
bullets(s, [
    ("FarmTwin Engine (open): GGA hydraulic solver, component/emitter library, FAO-56 agronomy — on GitHub.", 0, DARKTX, False),
    ("MVP: a runnable farm digital-twin simulator.", 0, DARKTX, False),
    ("Deep dossier: 13+ technical docs (solver math, sensors, IoT/fertigation, agronomy, optimisation, twin) + bibliography.", 0, DARKTX, False),
    ("Pilot land: 15 acres in the target panchayat.", 0, DARKTX, False),
    ("Brand: FarmTwin (trademark due-diligence done).", 0, DARKTX, False),
], y=Inches(1.5), size=20)

# --------------------------------------------------------- Slide 13: Ask
s = add_slide(); bg(s, WHITE); title_bar(s, "The ask — KSUM grant")
bullets(s, [
    ("Build the field-ready FarmTwin Runtime controller (solar, wireless) + sensor kit.", 0, DARKTX, False),
    ("Deploy a calibrated pilot on 15 acres + one Moolathara CMI zone.", 0, DARKTX, False),
    ("Validate with IIT Palakkad / KAU; publish a performance report.", 0, DARKTX, False),
    ("Use the pilot to win the first Moolathara CMI design / monitoring contract.", 0, DARKTX, False),
], y=Inches(1.6), size=21)
textbox(s, Inches(0.6), Inches(5.6), Inches(12.1), Inches(0.9),
        [[("12-month milestones: working controller → pilot live → validated "
           "water-saving & yield → first paid contract.", 16, GREY, True)]])

# --------------------------------------------------------- Slide 14: Vision
s = add_slide(); bg(s, INK)
rect(s, 0, 0, SW, Inches(0.12), WATER)
textbox(s, Inches(0.7), Inches(1.0), Inches(12), Inches(0.8),
        [[("Vision", 34, WHITE, True)]])
textbox(s, Inches(0.7), Inches(2.0), Inches(12), Inches(3.5),
        [[("A canal reached my village after 60 years. But concrete alone has failed "
           "us before.", 24, LIGHT, True)],
         [("FarmTwin makes sure that this time the water reaches the roots — and that "
           "we can prove it, drop by drop.", 24, LEAF, True)],
         [("We start with one canal in Chittur. We're building the system that lets any "
           "farmer, anywhere, grow more with less water.", 20, LIGHT, False)]])
textbox(s, Inches(0.7), Inches(5.7), Inches(12), Inches(1.0),
        [[("From simulating the world's factories — to securing our farmers' water.",
           22, WATER, True)],
         [("FarmTwin · github.com/arunthanga/FarmTwin", 14,
           RGBColor(0x9F,0xB3,0xC8), False)]], align=PP_ALIGN.CENTER)

import os
out = os.path.join(os.path.dirname(__file__), "FarmTwin-KSUM.pptx")
prs.save(out)
print("WROTE", out)
