#!/usr/bin/env python3
"""Build the FarmTwin pitch deck as both PPTX and PDF.

Single source of slide content (SLIDES) is rendered to:
  - FarmTwin-pitch.pptx  (python-pptx)
  - FarmTwin-pitch.pdf   (reportlab)

Run:
    pip install python-pptx reportlab
    python build_deck.py
"""

from pathlib import Path

# --- Brand palette ---------------------------------------------------------
GREEN = (0x1B, 0x7A, 0x43)   # FarmTwin green (agri)
BLUE = (0x1E, 0x73, 0xBE)    # water / tech accent
SLATE = (0x23, 0x2B, 0x33)   # body text
MUTED = (0x5A, 0x67, 0x72)   # secondary text
LIGHT = (0xF2, 0xF6, 0xF3)   # light panel

OUT_DIR = Path(__file__).resolve().parent

# --- Slide content ---------------------------------------------------------
# Each slide: dict with kind + fields. Bullets may be (text, level) tuples.

SLIDES = [
    {
        "kind": "title",
        "title": "FarmTwin",
        "subtitle": "Agri Digital-Twin & Precision-Operations Platform",
        "tagline": "Simulate every field. Irrigate by need, not by habit.",
        "footer": "Live pilot: 15-acre farm, Eruthempathy, Chittur (Palakkad) | KSUM / AgriNext",
    },
    {
        "kind": "content",
        "title": "The Problem",
        "subtitle": "Kerala's rain-shadow belt: water arrived, scheduling did not.",
        "bullets": [
            ("Chittur (Palakkad) gets ~1,000-1,250 mm rain vs the state's ~3,000 mm.", 0),
            ("The 2025 Moolathara canal extension now serves ~3,575 ha with drip/lift water.", 0),
            ("But farmers still irrigate by calendar and habit, not by crop-and-soil need:", 0),
            ("Over-irrigation wastes the scarce canal water just secured.", 1),
            ("Under-irrigation causes crop stress and yield loss.", 1),
            ("No plot-level, day-by-day decision tool exists for this microclimate.", 0),
            ("FPOs have no shared layer to budget water or plan harvests across 650+ ha.", 0),
        ],
    },
    {
        "kind": "metrics",
        "title": "Quantified Pain",
        "subtitle": "Why the bottleneck is now optimization, not availability.",
        "metrics": [
            ("30-50%", "of applied water wasted by flood/over-irrigation vs precise scheduling"),
            ("10-25%", "of potential yield lost to mistimed irrigation & nutrient stress"),
            ("~30%", "of manual irrigation labour reducible with scheduled precision"),
        ],
    },
    {
        "kind": "content",
        "title": "The Solution",
        "subtitle": "A farm digital twin built on a physics + agronomy engine.",
        "bullets": [
            ("Discretize a field into a computational mesh of zones - like a finite-element mesh.", 0),
            ("A time-stepped water-balance + crop-growth solver advances each zone, day by day.", 0),
            ("Ingests soil type, crop, emitter layout, weather/ET and low-cost soil-moisture data.", 0),
            ("Answers per zone, per day: irrigate or not, how many litres, and why.", 0),
            ("Optimizes the irrigation/fertigation schedule to minimize water & cost at target yield.", 0),
            ("Run what-if scenarios (flood vs drip vs optimized, heat wave) before acting in the field.", 0),
        ],
    },
    {
        "kind": "twocol",
        "title": "Why It's Defensible",
        "subtitle": "20 years of CAE / simulation, re-pointed at agriculture.",
        "left_head": "CAE / simulation skill",
        "right_head": "FarmTwin application",
        "rows": [
            ("Meshing algorithms", "Discretize the farm into adaptive irrigation/soil zones"),
            ("FE / CFD field solvers", "Water-balance + moisture & nutrient transport across the mesh"),
            ("Assembly-line simulation", "Time-stepped scheduling under a daily water budget"),
            ("DOE & optimization", "What-if scenarios + optimal irrigation/fertigation schedules"),
            ("Validation vs test data", "Calibrate the twin against sensors & observed yield"),
        ],
        "note": "The moat is the calibrated simulation engine - not another IoT dashboard.",
    },
    {
        "kind": "twocol",
        "title": "Two Products, One Shared Engine",
        "subtitle": "A calibration feedback loop that compounds over time.",
        "left_head": "FarmTwin Studio (pre-install)",
        "right_head": "FarmTwin Runtime (operations)",
        "rows": [
            ("Used once, before installation", "Runs continuously, after installation"),
            ("Survey -> simulate -> optimize", "Sense -> decide -> actuate"),
            ("Recommends top 2-3 designs, BoM", "Live control, twin state, alerts, yield"),
            ("Designer / agronomist / dealer", "Farmer / FPO operator (autonomous)"),
        ],
        "note": "Shared FarmTwin Engine (krishiflow): Runtime re-estimates parameters and feeds Studio - designs keep improving.",
    },
    {
        "kind": "content",
        "title": "Product & MVP",
        "subtitle": "A working prototype already demonstrates the core value.",
        "bullets": [
            ("Web app: a farm is a mesh of zones; each zone is modelled and solved daily.", 0),
            ("Scenario comparison - flood vs uniform drip vs twin-optimized.", 0),
            ("Outputs: water saved, stress-days avoided, projected yield.", 0),
            ("Hardware-light: cheap/existing soil & weather sensors + free weather/ET data.", 0),
            ("MVP runs in any browser (mvp/index.html) - no build step.", 0),
        ],
    },
    {
        "kind": "content",
        "title": "Market & Beachhead",
        "subtitle": "One FPO sale reaches hundreds of farmers.",
        "bullets": [
            ("Primary pilot: founder's 15-acre farm + 8-12 neighbouring vegetable/mango farmers.", 0),
            ("Beachhead buyer: FPOs / FPCs in the Chittur rain-shadow cluster (1,500+ farmers, 650 ha).", 0),
            ("Routed through the World Bank KERA / AgriNext program.", 1),
            ("Scale buyer: Krishi Bhavans & Dept. of Agriculture in other rain-shadow districts.", 0),
            ("Phase 2: lenders, insurers and input suppliers buying the data/risk layer.", 0),
        ],
    },
    {
        "kind": "content",
        "title": "Why Now",
        "subtitle": "Three tailwinds line up in 2025-2026.",
        "bullets": [
            ("Canal water arrived in 2025 - scheduling is suddenly the binding constraint.", 0),
            ("World Bank KERA / AgriNext (2026) is funding & distributing exactly this category to FPOs.", 0),
            ("Cheap soil/weather sensing + free ET data make a software-led, hardware-light product viable.", 0),
        ],
    },
    {
        "kind": "content",
        "title": "Pilot & Traction Plan",
        "subtitle": "Milestone-disbursed, verifiable success criteria.",
        "bullets": [
            ("M1 (month 2): sensors deployed on pilot farm; twin ingesting live data.", 0),
            ("M2 (month 4): solver calibrated to a defined soil-moisture error band.", 0),
            ("M3 (month 6): scenario engine shows >=30% water saving; 1 signed FPO pilot MoU.", 0),
            ("Target: >=30% simulated water saving vs flood baseline at equal yield.", 0),
        ],
    },
    {
        "kind": "content",
        "title": "Funding Roadmap",
        "subtitle": "Sequence non-dilutive money first; prove traction, then scale.",
        "bullets": [
            ("1. Idea Grant - up to Rs.3L (MVP + pilot plan).", 0),
            ("2. AgriNext / KERA pilot - up to Rs.25L (FPO match).", 0),
            ("3. Productisation Grant - up to Rs.7L (MVP -> product).", 0),
            ("4. Seed Fund / Seed Loan - up to Rs.15L (soft loan vs orders).", 0),
            ("5. Scale-up Seed Fund - up to Rs.25L (repeatable paid deployments).", 0),
            ("6. KERA Productive Alliance - up to Rs.2cr (FPC alliance, 60% matching).", 0),
        ],
    },
    {
        "kind": "twocol",
        "title": "Use of Funds & The Ask",
        "subtitle": "KSUM Idea Grant: Rs.3,00,000",
        "left_head": "Item",
        "right_head": "Amount (Rs.)",
        "rows": [
            ("Field instrumentation (sensors, gateway)", "90,000"),
            ("Cloud + dev tooling (12 months)", "60,000"),
            ("Solver/MVP hardening (contractor time)", "90,000"),
            ("Calibration & field validation", "45,000"),
            ("IP / legal (patent search, incorporation)", "15,000"),
        ],
        "note": "Ask: Rs.3 lakh Idea Grant + KSUM incubation + AgriNext deployment introductions.",
    },
    {
        "kind": "closing",
        "title": "FarmTwin",
        "subtitle": "I have spent 20 years meshing and simulating physical systems.",
        "tagline": "A farm is just another domain to mesh and solve - and I own the test rig.",
        "footer": "FarmTwin Studio  -  FarmTwin Runtime  -  FarmTwin Engine",
    },
]


# ---------------------------------------------------------------------------
# PPTX renderer
# ---------------------------------------------------------------------------
def build_pptx(path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    EMU = 914400
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    def rgb(c):
        return RGBColor(*c)

    def add_rect(slide, x, y, w, h, color, line=None):
        from pptx.enum.shapes import MSO_SHAPE
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(color)
        if line is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = rgb(line)
        shp.shadow.inherit = False
        return shp

    def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 space_after=6):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        for i, (text, size, color, bold) in enumerate(runs):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.space_after = Pt(space_after)
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = rgb(color)
            r.font.name = "Calibri"
        return tb

    def header(slide, title, subtitle):
        add_rect(slide, 0, 0, SW, Inches(1.45), GREEN)
        add_rect(slide, 0, Inches(1.45), SW, Emu(int(0.06 * EMU)), BLUE)
        add_text(slide, Inches(0.6), Inches(0.18), Inches(12), Inches(0.8),
                 [(title, 30, (255, 255, 255), True)], anchor=MSO_ANCHOR.MIDDLE,
                 space_after=0)
        if subtitle:
            add_text(slide, Inches(0.62), Inches(0.92), Inches(12), Inches(0.5),
                     [(subtitle, 14, (0xD8, 0xEC, 0xDF), False)], space_after=0)

    def render_title(s, blue=False):
        slide = prs.slides.add_slide(blank)
        bg = BLUE if blue else GREEN
        add_rect(slide, 0, 0, SW, SH, bg)
        add_rect(slide, 0, Inches(2.55), SW, Emu(int(0.07 * EMU)),
                 (255, 255, 255) if blue else BLUE)
        add_text(slide, Inches(0.8), Inches(1.55), Inches(11.7), Inches(1.1),
                 [(s["title"], 60, (255, 255, 255), True)], space_after=0)
        add_text(slide, Inches(0.82), Inches(2.75), Inches(11.7), Inches(0.8),
                 [(s["subtitle"], 22, (0xE6, 0xF0, 0xEA), False)], space_after=0)
        add_text(slide, Inches(0.82), Inches(3.7), Inches(11.7), Inches(0.8),
                 [(s["tagline"], 20, (255, 255, 255), True)], space_after=0)
        add_text(slide, Inches(0.82), Inches(6.7), Inches(11.7), Inches(0.6),
                 [(s["footer"], 13, (0xCF, 0xDE, 0xD5), False)], space_after=0)

    def render_content(s):
        slide = prs.slides.add_slide(blank)
        header(slide, s["title"], s.get("subtitle"))
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.85), Inches(12), Inches(5.3))
        tf = tb.text_frame
        tf.word_wrap = True
        first = True
        for text, level in s["bullets"]:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.level = level
            p.space_after = Pt(10)
            bullet = "-  " if level else "\u25B8  "
            r = p.add_run()
            r.text = bullet + text
            r.font.size = Pt(20 if level == 0 else 17)
            r.font.bold = False
            r.font.color.rgb = rgb(SLATE if level == 0 else MUTED)
            r.font.name = "Calibri"

    def render_metrics(s):
        slide = prs.slides.add_slide(blank)
        header(slide, s["title"], s.get("subtitle"))
        n = len(s["metrics"])
        gap = Inches(0.4)
        total_w = SW - Inches(1.2) - gap * (n - 1)
        cw = int(total_w / n)
        x = Inches(0.6)
        y = Inches(2.4)
        ch = Inches(3.3)
        for big, label in s["metrics"]:
            add_rect(slide, x, y, cw, ch, LIGHT)
            add_rect(slide, x, y, cw, Emu(int(0.08 * EMU)), GREEN)
            add_text(slide, x, y + Inches(0.7), cw, Inches(1.2),
                     [(big, 54, GREEN, True)], align=PP_ALIGN.CENTER, space_after=0)
            add_text(slide, x + Inches(0.25), y + Inches(1.9),
                     cw - Inches(0.5), Inches(1.3),
                     [(label, 16, SLATE, False)], align=PP_ALIGN.CENTER, space_after=0)
            x += cw + gap

    def render_twocol(s):
        slide = prs.slides.add_slide(blank)
        header(slide, s["title"], s.get("subtitle"))
        x0 = Inches(0.6)
        y0 = Inches(1.95)
        col_w = Inches(6.0)
        gap = Inches(0.13)
        right_x = x0 + col_w + gap
        rh = Inches(0.72)
        # headers
        add_rect(slide, x0, y0, col_w, Inches(0.5), GREEN)
        add_rect(slide, right_x, y0, col_w, Inches(0.5), BLUE)
        add_text(slide, x0 + Inches(0.15), y0 + Inches(0.03), col_w - Inches(0.3),
                 Inches(0.45), [(s["left_head"], 16, (255, 255, 255), True)],
                 anchor=MSO_ANCHOR.MIDDLE, space_after=0)
        add_text(slide, right_x + Inches(0.15), y0 + Inches(0.03),
                 col_w - Inches(0.3), Inches(0.45),
                 [(s["right_head"], 16, (255, 255, 255), True)],
                 anchor=MSO_ANCHOR.MIDDLE, space_after=0)
        y = y0 + Inches(0.5) + gap
        for i, (l, r) in enumerate(s["rows"]):
            band = LIGHT if i % 2 == 0 else (255, 255, 255)
            add_rect(slide, x0, y, col_w, rh, band)
            add_rect(slide, right_x, y, col_w, rh, band)
            add_text(slide, x0 + Inches(0.18), y, col_w - Inches(0.36), rh,
                     [(l, 15, SLATE, True)], anchor=MSO_ANCHOR.MIDDLE, space_after=0)
            add_text(slide, right_x + Inches(0.18), y, col_w - Inches(0.36), rh,
                     [(r, 15, SLATE, False)], anchor=MSO_ANCHOR.MIDDLE, space_after=0)
            y += rh + gap
        if s.get("note"):
            add_text(slide, x0, y + Inches(0.05), Inches(12.1), Inches(0.9),
                     [(s["note"], 15, GREEN, True)], space_after=0)

    for s in SLIDES:
        kind = s["kind"]
        if kind == "title":
            render_title(s, blue=False)
        elif kind == "closing":
            render_title(s, blue=True)
        elif kind == "content":
            render_content(s)
        elif kind == "metrics":
            render_metrics(s)
        elif kind == "twocol":
            render_twocol(s)

    prs.save(str(path))
    print(f"wrote {path}")


# ---------------------------------------------------------------------------
# PDF renderer
# ---------------------------------------------------------------------------
def build_pdf(path: Path) -> None:
    from reportlab.lib.pagesizes import landscape
    from reportlab.pdfgen import canvas
    from reportlab.lib.utils import simpleSplit

    PW, PH = (960, 540)  # 16:9 points
    c = canvas.Canvas(str(path), pagesize=(PW, PH))

    def col(rgbt):
        return tuple(v / 255 for v in rgbt)

    def set_fill(rgbt):
        c.setFillColorRGB(*col(rgbt))

    def rect(x, y, w, h, rgbt):
        set_fill(rgbt)
        c.rect(x, y, w, h, stroke=0, fill=1)

    def wrap(text, font, size, max_w):
        return simpleSplit(text, font, size, max_w)

    def draw_wrapped(text, x, y, font, size, rgbt, max_w, leading, align="left"):
        set_fill(rgbt)
        c.setFont(font, size)
        lines = wrap(text, font, size, max_w)
        for line in lines:
            if align == "center":
                c.drawCentredString(x, y, line)
            else:
                c.drawString(x, y, line)
            y -= leading
        return y

    def header(title, subtitle):
        rect(0, PH - 95, PW, 95, GREEN)
        rect(0, PH - 99, PW, 4, BLUE)
        set_fill((255, 255, 255))
        c.setFont("Helvetica-Bold", 26)
        c.drawString(45, PH - 52, title)
        if subtitle:
            set_fill((0xD8, 0xEC, 0xDF))
            c.setFont("Helvetica", 13)
            c.drawString(46, PH - 80, subtitle)

    def title_page(s, blue=False):
        bg = BLUE if blue else GREEN
        rect(0, 0, PW, PH, bg)
        rect(0, PH - 215, PW, 5, (255, 255, 255) if blue else BLUE)
        set_fill((255, 255, 255))
        c.setFont("Helvetica-Bold", 54)
        c.drawString(60, PH - 195, s["title"])
        set_fill((0xE6, 0xF0, 0xEA))
        c.setFont("Helvetica", 20)
        draw_wrapped(s["subtitle"], 62, PH - 250, "Helvetica", 20,
                     (0xE6, 0xF0, 0xEA), PW - 120, 26)
        set_fill((255, 255, 255))
        c.setFont("Helvetica-Bold", 18)
        draw_wrapped(s["tagline"], 62, PH - 320, "Helvetica-Bold", 18,
                     (255, 255, 255), PW - 120, 24)
        set_fill((0xCF, 0xDE, 0xD5))
        c.setFont("Helvetica", 11)
        draw_wrapped(s["footer"], 62, 55, "Helvetica", 11,
                     (0xCF, 0xDE, 0xD5), PW - 120, 16)
        c.showPage()

    def content_page(s):
        header(s["title"], s.get("subtitle"))
        y = PH - 135
        for text, level in s["bullets"]:
            if level == 0:
                bullet = "\u25B8  "
                font, size, rgbt, indent = "Helvetica", 15, SLATE, 50
            else:
                bullet = "-  "
                font, size, rgbt, indent = "Helvetica", 13, MUTED, 80
            y = draw_wrapped(bullet + text, indent, y, font, size, rgbt,
                             PW - indent - 60, 20)
            y -= 8
        c.showPage()

    def metrics_page(s):
        header(s["title"], s.get("subtitle"))
        m = s["metrics"]
        n = len(m)
        gap = 25
        margin = 45
        cw = (PW - 2 * margin - gap * (n - 1)) / n
        ch = 230
        y = PH - 380
        x = margin
        for big, label in m:
            rect(x, y, cw, ch, LIGHT)
            rect(x, y + ch - 6, cw, 6, GREEN)
            set_fill(GREEN)
            c.setFont("Helvetica-Bold", 46)
            c.drawCentredString(x + cw / 2, y + ch - 80, big)
            set_fill(SLATE)
            c.setFont("Helvetica", 13)
            lines = wrap(label, "Helvetica", 13, cw - 30)
            ly = y + ch - 130
            for line in lines:
                c.drawCentredString(x + cw / 2, ly, line)
                ly -= 17
            x += cw + gap
        c.showPage()

    def twocol_page(s):
        header(s["title"], s.get("subtitle"))
        margin = 45
        gap = 10
        col_w = (PW - 2 * margin - gap) / 2
        x_l = margin
        x_r = margin + col_w + gap
        y = PH - 135
        head_h = 32
        rect(x_l, y - head_h, col_w, head_h, GREEN)
        rect(x_r, y - head_h, col_w, head_h, BLUE)
        set_fill((255, 255, 255))
        c.setFont("Helvetica-Bold", 13)
        c.drawString(x_l + 12, y - head_h + 11, s["left_head"])
        c.drawString(x_r + 12, y - head_h + 11, s["right_head"])
        y -= head_h
        for i, (l, r) in enumerate(s["rows"]):
            row_h = 46
            band = LIGHT if i % 2 == 0 else (255, 255, 255)
            rect(x_l, y - row_h, col_w, row_h, band)
            rect(x_r, y - row_h, col_w, row_h, band)
            set_fill(SLATE)
            c.setFont("Helvetica-Bold", 12)
            ll = wrap(l, "Helvetica-Bold", 12, col_w - 24)
            ty = y - 18
            for line in ll:
                c.drawString(x_l + 12, ty, line)
                ty -= 15
            c.setFont("Helvetica", 12)
            rl = wrap(r, "Helvetica", 12, col_w - 24)
            ty = y - 18
            for line in rl:
                c.drawString(x_r + 12, ty, line)
                ty -= 15
            y -= row_h + 4
        if s.get("note"):
            draw_wrapped(s["note"], margin, y - 18, "Helvetica-Bold", 12,
                         GREEN, PW - 2 * margin, 16)
        c.showPage()

    for s in SLIDES:
        kind = s["kind"]
        if kind == "title":
            title_page(s, blue=False)
        elif kind == "closing":
            title_page(s, blue=True)
        elif kind == "content":
            content_page(s)
        elif kind == "metrics":
            metrics_page(s)
        elif kind == "twocol":
            twocol_page(s)

    c.save()
    print(f"wrote {path}")


if __name__ == "__main__":
    build_pptx(OUT_DIR / "FarmTwin-pitch.pptx")
    build_pdf(OUT_DIR / "FarmTwin-pitch.pdf")
